import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    report_dir = os.path.join(os.path.dirname(__file__), '..', 'report')
    assets_dir = os.path.join(report_dir, 'assets')
    output_pptx = os.path.join(report_dir, 'PennyWise_Project_Presentation.pptx')

    print(f"[PPTX] Creating PowerPoint presentation: {output_pptx}")
    prs = Presentation()

    # Set 16:9 Widescreen dimensions (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    BG_DARK = RGBColor(15, 23, 42)       # #0F172A (Slate 900)
    CARD_BG = RGBColor(30, 41, 59)      # #1E293B (Slate 800)
    TEXT_WHITE = RGBColor(248, 250, 252) # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)# #94A3B8
    PRIMARY = RGBColor(99, 102, 241)    # #6366F1 (Indigo)
    ACCENT_GREEN = RGBColor(16, 185, 129)# #10B981 (Emerald)
    ACCENT_AMBER = RGBColor(245, 158, 11)# #F59E0B (Amber)

    blank_layout = prs.slide_layouts[6]

    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_header(slide, title_text, category_text="PENNYWISE PROJECT PRESENTATION"):
        # Header category badge
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        # Title
        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf2 = tx_box2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)

    # Decorative Card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = PRIMARY
    card.line.width = Pt(2)

    tb = slide1.shapes.add_textbox(Inches(2.0), Inches(1.8), Inches(9.333), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "PennyWise"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Personal Finance & Behavioral Economics Platform"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.space_after = Pt(30)

    p3 = tf.add_paragraph()
    p3.text = "Presented by: Sneha  |  Roll No: 2105152004"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_GREEN

    p4 = tf.add_paragraph()
    p4.text = "Bachelor of Technology in Computer Science & Engineering\nSchool of Computer Engineering — KIIT Deemed to be University"
    p4.font.size = Pt(14)
    p4.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 2: Problem Statement & Motivation
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_header(slide2, "Problem Statement & Project Motivation")

    # Left Box: Problem
    card_l = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = RGBColor(239, 68, 68)

    tb_l = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "Traditional Finance Apps"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)

    points_l = [
        "Record numbers without explaining behavioral causes.",
        "Ignore emotional spending triggers (stress, boredom, impulse).",
        "Fail to quantify long-term impact on savings goals.",
        "Lack interactive feedback to prevent impulse spending."
    ]
    for pt in points_l:
        p_item = tf_l.add_paragraph()
        p_item.text = "• " + pt
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = TEXT_WHITE
        p_item.space_before = Pt(12)

    # Right Box: Solution
    card_r = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = ACCENT_GREEN

    tb_r = slide2.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "PennyWise Solution"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    points_r = [
        "Combines transaction records with behavioral psychology.",
        "Transparent 6-weight Regret Predictor Engine (0-100 score).",
        "Impulse Lock cooling-off timer with real-time AI advice.",
        "Gamification (XP, levels, badges, daily savings missions)."
    ]
    for pt in points_r:
        p_item = tf_r.add_paragraph()
        p_item.text = "• " + pt
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = TEXT_WHITE
        p_item.space_before = Pt(12)

    # ==========================================
    # SLIDE 3: System Architecture
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_header(slide3, "System Architecture & Tech Stack")

    img3 = os.path.join(assets_dir, 'fig_3_1_architecture.png')
    if os.path.exists(img3):
        slide3.shapes.add_picture(img3, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))

    # ==========================================
    # SLIDE 4: System Use Cases & Data Flow
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_header(slide4, "Use Case Model & Data Flow Architecture")

    img4_1 = os.path.join(assets_dir, 'fig_3_2_usecase.png')
    img4_2 = os.path.join(assets_dir, 'fig_3_4_dataflow.png')
    if os.path.exists(img4_1):
        slide4.shapes.add_picture(img4_1, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    if os.path.exists(img4_2):
        slide4.shapes.add_picture(img4_2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))

    # ==========================================
    # SLIDE 5: Financial Command Center (Dashboard)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_header(slide5, "Financial Command Center (Dashboard)")

    img5 = os.path.join(assets_dir, 'fig_4_1_dashboard.png')
    if os.path.exists(img5):
        slide5.shapes.add_picture(img5, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))

    # ==========================================
    # SLIDE 6: AI Behavior Suite
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)
    add_header(slide6, "AI Behavior Suite & Decision Tools")

    img6 = os.path.join(assets_dir, 'fig_4_5_ai_suite.png')
    if os.path.exists(img6):
        slide6.shapes.add_picture(img6, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))

    # ==========================================
    # SLIDE 7: Regret Predictor Engine
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7)
    add_header(slide7, "3-Panel Regret Predictor Engine")

    img7 = os.path.join(assets_dir, 'fig_4_7_regret_predictor.png')
    if os.path.exists(img7):
        slide7.shapes.add_picture(img7, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))

    # ==========================================
    # SLIDE 8: Dream Savings Goals
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8)
    add_header(slide8, "Dream Savings Goals & Progression")

    img8 = os.path.join(assets_dir, 'fig_4_8_dream_goals.png')
    if os.path.exists(img8):
        slide8.shapes.add_picture(img8, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))

    # ==========================================
    # SLIDE 9: Subscription Killer
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9)
    add_header(slide9, "Subscription Killer & Usage Analytics")

    img9 = os.path.join(assets_dir, 'fig_4_9_subscriptions.png')
    if os.path.exists(img9):
        slide9.shapes.add_picture(img9, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))

    # ==========================================
    # SLIDE 10: Gamification & Calendar
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10)
    add_header(slide10, "Gamification, XP & No-Spend Calendar")

    img10_1 = os.path.join(assets_dir, 'fig_4_10_gamification.png')
    img10_2 = os.path.join(assets_dir, 'fig_4_11_nospend_calendar.png')
    if os.path.exists(img10_1):
        slide10.shapes.add_picture(img10_1, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    if os.path.exists(img10_2):
        slide10.shapes.add_picture(img10_2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))

    # ==========================================
    # SLIDE 11: Testing & Quality Assurance
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide11)
    add_header(slide11, "Testing & Quality Assurance")

    # Card 1: 100% Pass Rate Box
    c1 = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = ACCENT_GREEN
    c1.line.width = Pt(2)

    tb1 = slide11.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "27 / 27 Automated Tests Passed"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    test_items = [
        "Authentication Tests (8/8 Passed) — bcrypt hashing, JWT validation, token isolation.",
        "Expenses CRUD Tests (5/5 Passed) — user-scoped expense operations & satisfaction rating.",
        "Budgeting & Reports Tests (5/5 Passed) — category budget thresholds & PDF output.",
        "AI Behavior Suite Tests (6/6 Passed) — Regret Predictor score calculations.",
        "System E2E Workflows (3/3 Passed) — full user lifecycle and account cleanup."
    ]
    for item in test_items:
        p_item = tf1.add_paragraph()
        p_item.text = "• " + item
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = TEXT_WHITE
        p_item.space_before = Pt(10)

    # Card 2: Security & Isolation
    c2 = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = PRIMARY
    c2.line.width = Pt(2)

    tb2 = slide11.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Security & Quality Standards"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    sec_items = [
        "JWT Authentication Guarding all REST API routes.",
        "Bcrypt Password Hashing with multi-round salt.",
        "Strict User Isolation via parameterized SQL query filters.",
        "Cascading Foreign Key Deletion to prevent orphaned data.",
        "Environment Secret Isolation (.env configuration)."
    ]
    for item in sec_items:
        p_item = tf2.add_paragraph()
        p_item.text = "• " + item
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = TEXT_WHITE
        p_item.space_before = Pt(14)

    # ==========================================
    # SLIDE 12: Conclusion & Q&A
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide12)

    card12 = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    card12.fill.solid()
    card12.fill.fore_color.rgb = CARD_BG
    card12.line.color.rgb = ACCENT_GREEN
    card12.line.width = Pt(2)

    tb12 = slide12.shapes.add_textbox(Inches(2.0), Inches(1.8), Inches(9.333), Inches(3.8))
    tf12 = tb12.text_frame
    tf12.word_wrap = True
    tf12.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf12.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    p2 = tf12.add_paragraph()
    p2.text = "PennyWise: AI-Powered Personal Finance Platform"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(16)

    p3 = tf12.add_paragraph()
    p3.text = "Questions & Answers"
    p3.font.size = Pt(26)
    p3.font.bold = True
    p3.font.color.rgb = PRIMARY
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)

    print(f"[PPTX] Saving PowerPoint presentation to: {output_pptx}")
    prs.save(output_pptx)
    print("[PPTX] PowerPoint Presentation Created Successfully!")

if __name__ == '__main__':
    build_presentation()
