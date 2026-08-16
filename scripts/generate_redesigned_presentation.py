import os
import shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_redesigned_presentation():
    report_dir = os.path.join(os.path.dirname(__file__), '..', 'report')
    assets_dir = os.path.join(report_dir, 'assets')
    output_pptx = os.path.join(report_dir, 'PennyWise_Redesigned_Presentation.pptx')
    target_pptx = os.path.join(report_dir, 'PennyWise_Project_Presentation.pptx')

    print(f"[PPTX] Generating redesigned 12-slide presentation: {output_pptx}")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme Colors
    BG_DARK = RGBColor(11, 15, 25)          # #0B0F19 (Deep Navy/Near Black)
    CARD_BG = RGBColor(24, 32, 47)         # #18202F (Dark Glassmorphism Card)
    CARD_BORDER = RGBColor(99, 102, 241)   # #6366F1 (Electric Indigo)
    TEXT_WHITE = RGBColor(248, 250, 252)    # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8
    TEXT_SUBTLE = RGBColor(203, 213, 225)  # #CBD5E1
    PRIMARY_PURPLE = RGBColor(129, 140, 248)# #818CF8
    MINT_GREEN = RGBColor(52, 211, 153)    # #34D399 (Positive)
    AMBER_WARN = RGBColor(251, 191, 36)    # #FBBF24 (Warning)
    RED_RISK = RGBColor(248, 113, 113)     # #F87171 (Risk)
    CYAN_BLUE = RGBColor(56, 189, 248)     # #38BDF8

    blank_layout = prs.slide_layouts[6]

    def apply_background(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK

    def add_top_brand(slide, slide_num, title_text, tagline="AI-POWERED PERSONAL FINANCE & BEHAVIORAL ECONOMICS"):
        # Header Badge
        tb_brand = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(10.0), Inches(0.35))
        tf_b = tb_brand.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.text = f"PENNYWISE  •  {tagline}"
        p_b.font.size = Pt(9)
        p_b.font.bold = True
        p_b.font.color.rgb = PRIMARY_PURPLE

        # Slide Title
        tb_title = slide.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(10.0), Inches(0.65))
        tf_t = tb_title.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

        # Slide Number Badge
        tb_num = slide.shapes.add_textbox(Inches(11.8), Inches(0.35), Inches(0.9), Inches(0.35))
        tf_n = tb_num.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = f"{slide_num:02d} / 12"
        p_n.alignment = PP_ALIGN.RIGHT
        p_n.font.size = Pt(10)
        p_n.font.bold = True
        p_n.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, border_color=CARD_BORDER, bg_color=CARD_BG):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (Striking Hero Dashboard Visual)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1)

    # Left Information Column
    tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(5.8), Inches(5.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "PennyWise"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_PURPLE
    p.space_after = Pt(4)

    p_sub = tf1.add_paragraph()
    p_sub.text = "Spend smarter. Think better. Save more."
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = MINT_GREEN
    p_sub.space_after = Pt(14)

    p_desc = tf1.add_paragraph()
    p_desc.text = "AI-Powered Personal Finance &\nBehavioral Economics Platform"
    p_desc.font.size = Pt(22)
    p_desc.font.bold = True
    p_desc.font.color.rgb = TEXT_WHITE
    p_desc.space_after = Pt(36)

    # Presenter Card
    add_card(slide1, 0.8, 4.6, 5.4, 2.1, border_color=PRIMARY_PURPLE)
    tb_pres = slide1.shapes.add_textbox(Inches(1.0), Inches(4.75), Inches(5.0), Inches(1.8))
    tf_p = tb_pres.text_frame
    tf_p.word_wrap = True

    p_by = tf_p.paragraphs[0]
    p_by.text = "Presented by:"
    p_by.font.size = Pt(11)
    p_by.font.color.rgb = TEXT_MUTED

    p_name = tf_p.add_paragraph()
    p_name.text = "Sneha  |  Roll No: 2305495"
    p_name.font.size = Pt(18)
    p_name.font.bold = True
    p_name.font.color.rgb = MINT_GREEN
    p_name.space_after = Pt(6)

    p_univ = tf_p.add_paragraph()
    p_univ.text = "B.Tech Computer Science & Engineering\nSchool of Computer Engineering\nKIIT Deemed to be University"
    p_univ.font.size = Pt(12)
    p_univ.font.color.rgb = TEXT_SUBTLE

    # Right Hero Dashboard Screenshot
    dash_img = os.path.join(assets_dir, 'fig_4_1_dashboard.png')
    if os.path.exists(dash_img):
        add_card(slide1, 6.5, 0.9, 6.2, 5.8, border_color=CYAN_BLUE)
        slide1.shapes.add_picture(dash_img, Inches(6.6), Inches(1.0), Inches(6.0), Inches(5.6))

    # =========================================================================
    # SLIDE 2: PROBLEM → SOLUTION (Visual Transition)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2)
    add_top_brand(slide2, 2, "The Problem vs. PennyWise Solution")

    # Left Box: Problem
    add_card(slide2, 0.8, 1.5, 5.2, 5.3, border_color=RED_RISK)
    tb_p = slide2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(4.8), Inches(4.9))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True

    p = tf_p.paragraphs[0]
    p.text = "THE PROBLEM"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED_RISK

    prob_items = [
        ("Static Numbers Only", "Traditional finance apps record transactions without explaining underlying behavioral causes."),
        ("Ignored Emotional Triggers", "Impulse purchases driven by stress, boredom, or peer pressure are completely overlooked."),
        ("Hidden Goal Impact", "Users cannot visualize how small daily expenses delay long-term savings goals."),
        ("No Active Intervention", "No real-time cooling off or decision support before an impulse transaction occurs.")
    ]
    for title, desc in prob_items:
        p_t = tf_p.add_paragraph()
        p_t.text = "• " + title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.space_before = Pt(10)
        p_d = tf_p.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED

    # Arrow Divider
    arr = slide2.shapes.add_textbox(Inches(6.05), Inches(3.6), Inches(1.2), Inches(1.0))
    tf_a = arr.text_frame
    p_a = tf_a.paragraphs[0]
    p_a.text = "➔"
    p_a.alignment = PP_ALIGN.CENTER
    p_a.font.size = Pt(36)
    p_a.font.bold = True
    p_a.font.color.rgb = PRIMARY_PURPLE

    # Right Box: Solution
    add_card(slide2, 7.3, 1.5, 5.2, 5.3, border_color=MINT_GREEN)
    tb_s = slide2.shapes.add_textbox(Inches(7.5), Inches(1.7), Inches(4.8), Inches(4.9))
    tf_s = tb_s.text_frame
    tf_s.word_wrap = True

    p = tf_s.paragraphs[0]
    p.text = "PENNYWISE SOLUTION"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = MINT_GREEN

    sol_items = [
        ("Behavioral Finance + AI", "Combines transaction records with psychological mood tracking & financial health scoring."),
        ("Transparent Regret Predictor", "Computes a 0–100 risk score based on 6 weighted financial & behavioral components."),
        ("Impulse Lock & AI Advisor", "Enforces a 30-second cooling-off timer with real-time affordability analysis."),
        ("Gamified Progression", "Drives savings discipline through XP, levels, badges, and daily savings missions.")
    ]
    for title, desc in sol_items:
        p_t = tf_s.add_paragraph()
        p_t.text = "✓ " + title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = MINT_GREEN
        p_t.space_before = Pt(10)
        p_d = tf_s.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_SUBTLE

    # =========================================================================
    # SLIDE 3: WHAT IS PENNYWISE? / SYSTEM ARCHITECTURE Pipeline
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3)
    add_top_brand(slide3, 3, "System Architecture & Engineering Pipeline")

    arch_img = os.path.join(assets_dir, 'fig_3_1_architecture.png')
    if os.path.exists(arch_img):
        add_card(slide3, 0.6, 1.4, 7.8, 5.5, border_color=PRIMARY_PURPLE)
        slide3.shapes.add_picture(arch_img, Inches(0.7), Inches(1.5), Inches(7.6), Inches(5.3))

    # Right Architecture Callouts
    add_card(slide3, 8.6, 1.4, 4.1, 5.5, border_color=CYAN_BLUE)
    tb_ac = slide3.shapes.add_textbox(Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.1))
    tf_ac = tb_ac.text_frame
    tf_ac.word_wrap = True

    p = tf_ac.paragraphs[0]
    p.text = "Tech Stack Summary"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN_BLUE

    tech_stack = [
        ("Presentation Layer", "HTML5, CSS3 Glassmorphic UI, Vanilla JavaScript SPA Controller, Chart.js Visualizations."),
        ("Application Backend", "Node.js REST API, Express.js Middleware Stack, JWT Auth Guard, bcrypt Hashing."),
        ("Behavioral AI Engine", "6-Weight Regret Risk Engine, Should I Buy It Advisor, Opportunity Cost Simulator."),
        ("Persistence & Reports", "SQLite Relational Database (finpilot.db, 21 Tables), PDFKit Report Generator.")
    ]
    for layer, desc in tech_stack:
        pt = tf_ac.add_paragraph()
        pt.text = "• " + layer
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_WHITE
        pt.space_before = Pt(10)
        pd = tf_ac.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: USE CASE + DATA FLOW Workflow
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4)
    add_top_brand(slide4, 4, "Use Case Model & Data-Flow Workflow")

    uc_img = os.path.join(assets_dir, 'fig_3_2_usecase.png')
    df_img = os.path.join(assets_dir, 'fig_3_4_dataflow.png')

    if os.path.exists(uc_img):
        add_card(slide4, 0.6, 1.4, 5.8, 5.5, border_color=PRIMARY_PURPLE)
        slide4.shapes.add_picture(uc_img, Inches(0.7), Inches(1.5), Inches(5.6), Inches(5.3))

    if os.path.exists(df_img):
        add_card(slide4, 6.7, 1.4, 6.0, 5.5, border_color=MINT_GREEN)
        slide4.shapes.add_picture(df_img, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3))

    # =========================================================================
    # SLIDE 5: FINANCIAL COMMAND CENTER (Dashboard Hero)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5)
    add_top_brand(slide5, 5, "Financial Command Center & Dashboard")

    if os.path.exists(dash_img):
        add_card(slide5, 0.6, 1.4, 8.2, 5.5, border_color=CYAN_BLUE)
        slide5.shapes.add_picture(dash_img, Inches(0.7), Inches(1.5), Inches(8.0), Inches(5.3))

    # Right Stat Callout Badges
    add_card(slide5, 9.0, 1.4, 3.7, 5.5, border_color=PRIMARY_PURPLE)
    tb_d = slide5.shapes.add_textbox(Inches(9.15), Inches(1.6), Inches(3.4), Inches(5.1))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True

    p = tf_d.paragraphs[0]
    p.text = "Key Indicators"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_PURPLE

    dash_metrics = [
        ("Financial Health Score", "Calculated 0-100 score from savings rate, impulse frequency & mood stability."),
        ("Take-Home Income", "Calculates net balance after fixed payroll & tax deductions."),
        ("Monthly Spending Trend", "Compares current vs. previous month expenses by category."),
        ("Impulse & Mood Tracker", "Maps spending to emotional states (Happy, Stressed, Bored).")
    ]
    for label, desc in dash_metrics:
        pt = tf_d.add_paragraph()
        pt.text = "• " + label
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = MINT_GREEN
        pt.space_before = Pt(10)
        pd = tf_d.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 6: AI BEHAVIOR SUITE Ecosystem
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6)
    add_top_brand(slide6, 6, "AI Behavior Suite & Decision Tools")

    ai_img = os.path.join(assets_dir, 'fig_4_5_ai_suite.png')
    if os.path.exists(ai_img):
        add_card(slide6, 0.6, 1.4, 8.2, 5.5, border_color=PRIMARY_PURPLE)
        slide6.shapes.add_picture(ai_img, Inches(0.7), Inches(1.5), Inches(8.0), Inches(5.3))

    # Right Tool List
    add_card(slide6, 9.0, 1.4, 3.7, 5.5, border_color=AMBER_WARN)
    tb_ai = slide6.shapes.add_textbox(Inches(9.15), Inches(1.6), Inches(3.4), Inches(5.1))
    tf_ai = tb_ai.text_frame
    tf_ai.word_wrap = True

    p = tf_ai.paragraphs[0]
    p.text = "8 AI Decision Tools"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AMBER_WARN

    ai_tools = [
        ("Should I Buy It? Advisor", "Evaluates affordability against active savings goals."),
        ("Regret Predictor Engine", "6-weighted risk score with goal delay estimates."),
        ("Opportunity Cost Tools", "Translates price into daily coffee/meal equivalents."),
        ("Future Savings Simulator", "Models 1 & 5-year savings from category cutbacks."),
        ("Hidden Expense Detector", "Identifies micro-spending accumulation."),
        ("Weekly Spending Roast", "Generates humorous behavioral observations.")
    ]
    for label, desc in ai_tools:
        pt = tf_ai.add_paragraph()
        pt.text = "• " + label
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_WHITE
        pt.space_before = Pt(6)
        pd = tf_ai.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 7: REGRET PREDICTOR — HOW IT WORKS (3-Panel Hero)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7)
    add_top_brand(slide7, 7, "Regret Predictor Engine — 6 Weighted Component Analysis")

    rp_img = os.path.join(assets_dir, 'fig_4_7_regret_predictor.png')
    if os.path.exists(rp_img):
        add_card(slide7, 0.6, 1.4, 8.2, 5.5, border_color=RED_RISK)
        slide7.shapes.add_picture(rp_img, Inches(0.7), Inches(1.5), Inches(8.0), Inches(5.3))

    # Right Risk Weights Grid
    add_card(slide7, 9.0, 1.4, 3.7, 5.5, border_color=RED_RISK)
    tb_rp = slide7.shapes.add_textbox(Inches(9.15), Inches(1.6), Inches(3.4), Inches(5.1))
    tf_rp = tb_rp.text_frame
    tf_rp.word_wrap = True

    p = tf_rp.paragraphs[0]
    p.text = "Weighted Risk Model"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED_RISK

    weights = [
        ("Income Risk", "25%", "Price vs. monthly take-home pay."),
        ("Category Risk", "20%", "Historical category regret rates."),
        ("Mood Risk", "15%", "Emotional state (Stressed/Bored)."),
        ("Impulse Risk", "15%", "Unplanned purchase flag."),
        ("Goal Delay Risk", "15%", "Impact on active Dream Goals."),
        ("Purchase History", "10%", "Prior satisfaction feedback score.")
    ]
    for name, w, desc in weights:
        pt = tf_rp.add_paragraph()
        pt.text = f"• {name} ({w})"
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = AMBER_WARN
        pt.space_before = Pt(6)
        pd = tf_rp.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 8: DREAM SAVINGS GOALS (Hero)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_background(slide8)
    add_top_brand(slide8, 8, "Dream Savings Goals & Level Progression")

    dg_img = os.path.join(assets_dir, 'fig_4_8_dream_goals.png')
    if os.path.exists(dg_img):
        add_card(slide8, 0.6, 1.4, 8.2, 5.5, border_color=MINT_GREEN)
        slide8.shapes.add_picture(dg_img, Inches(0.7), Inches(1.5), Inches(8.0), Inches(5.3))

    # Right Progression Badges
    add_card(slide8, 9.0, 1.4, 3.7, 5.5, border_color=PRIMARY_PURPLE)
    tb_dg = slide8.shapes.add_textbox(Inches(9.15), Inches(1.6), Inches(3.4), Inches(5.1))
    tf_dg = tb_dg.text_frame
    tf_dg.word_wrap = True

    p = tf_dg.paragraphs[0]
    p.text = "Goal Progression Levels"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = MINT_GREEN

    levels = [
        ("Penny Beginner", "0% - 19% saved toward target."),
        ("Budget Explorer", "20% - 39% saved toward target."),
        ("Treasure Hunter", "40% - 59% saved toward target."),
        ("Savings Knight", "60% - 79% saved toward target."),
        ("Wealth Wizard", "80% - 99% saved toward target."),
        ("Financial Legend", "100% Goal Reached & Unlocked!")
    ]
    for name, desc in levels:
        pt = tf_dg.add_paragraph()
        pt.text = "• " + name
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_WHITE
        pt.space_before = Pt(6)
        pd = tf_dg.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(10)
        pd.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 9: SUBSCRIPTION KILLER + GAMIFICATION
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_background(slide9)
    add_top_brand(slide9, 9, "Subscription Killer & Gamification Engine")

    sub_img = os.path.join(assets_dir, 'fig_4_9_subscriptions.png')
    gam_img = os.path.join(assets_dir, 'fig_4_10_gamification.png')

    if os.path.exists(sub_img):
        add_card(slide9, 0.6, 1.4, 5.8, 5.5, border_color=AMBER_WARN)
        slide9.shapes.add_picture(sub_img, Inches(0.7), Inches(1.5), Inches(5.6), Inches(5.3))

    if os.path.exists(gam_img):
        add_card(slide9, 6.7, 1.4, 6.0, 5.5, border_color=PRIMARY_PURPLE)
        slide9.shapes.add_picture(gam_img, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3))

    # =========================================================================
    # SLIDE 10: PENNYWISE IN ACTION (Decision Flow Story)
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    apply_background(slide10)
    add_top_brand(slide10, 10, "PennyWise in Action — End-to-End Decision Flow")

    # Step Cards Flow
    steps = [
        ("1. User Intent", "User encounters an unplanned item (e.g. $250 Sneakers).", PRIMARY_PURPLE),
        ("2. Cooling-Off Lock", "Impulse Lock activates a 30-second cooling-off timer.", AMBER_WARN),
        ("3. AI Risk Analysis", "Regret Predictor evaluates 6 risk weights & affordability.", RED_RISK),
        ("4. Decision & Savings", "User reviews advice, cancels purchase, and earns XP!", MINT_GREEN)
    ]
    card_width = 2.7
    for idx, (stitle, sdesc, scolor) in enumerate(steps):
        left_pos = 0.6 + (idx * 3.1)
        add_card(slide10, left_pos, 1.6, card_width, 2.2, border_color=scolor)
        tb_s = slide10.shapes.add_textbox(Inches(left_pos + 0.15), Inches(1.75), Inches(card_width - 0.3), Inches(1.9))
        tf_s = tb_s.text_frame
        tf_s.word_wrap = True
        p = tf_s.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = scolor
        p_d = tf_s.add_paragraph()
        p_d.text = sdesc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_WHITE
        p_d.space_before = Pt(8)

    # Bottom Screenshot Preview
    buy_img = os.path.join(assets_dir, 'fig_4_6_buy_advisor.png')
    if os.path.exists(buy_img):
        add_card(slide10, 0.6, 4.1, 12.1, 2.8, border_color=CYAN_BLUE)
        slide10.shapes.add_picture(buy_img, Inches(0.7), Inches(4.2), Inches(11.9), Inches(2.6))

    # =========================================================================
    # SLIDE 11: TESTING & QUALITY ASSURANCE (100% Pass Rate Hero)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    apply_background(slide11)
    add_top_brand(slide11, 11, "Testing & Quality Assurance — 100% Verification")

    # Left Hero Stat Card
    add_card(slide11, 0.6, 1.4, 5.8, 5.5, border_color=MINT_GREEN)
    tb_t1 = slide11.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.4), Inches(5.1))
    tf_t1 = tb_t1.text_frame
    tf_t1.word_wrap = True

    p = tf_t1.paragraphs[0]
    p.text = "27 / 27 PASS"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = MINT_GREEN

    p_sub = tf_t1.add_paragraph()
    p_sub.text = "100% Automated Test Pass Rate across 5 Jest Test Suites"
    p_sub.font.size = Pt(13)
    p_sub.font.bold = True
    p_sub.font.color.rgb = TEXT_WHITE
    p_sub.space_after = Pt(16)

    t_suites = [
        ("Authentication Unit Tests", "8 / 8 Passed", "bcrypt hashing, JWT token validation & security isolation."),
        ("Expenses CRUD Unit Tests", "5 / 5 Passed", "User-scoped expense management & satisfaction ratings."),
        ("Budgeting & Reports Tests", "5 / 5 Passed", "Category budget limits, threshold alerts & PDF output."),
        ("AI Behavior Suite Tests", "6 / 6 Passed", "Regret Predictor scoring, advisor & simulator math."),
        ("System E2E Workflows", "3 / 3 Passed", "End-to-end user lifecycle and complete account cleanup.")
    ]
    for sname, spass, sdesc in t_suites:
        pt = tf_t1.add_paragraph()
        pt.text = f"✓ {sname} ({spass})"
        pt.font.size = Pt(12)
        pt.font.bold = True
        pt.font.color.rgb = MINT_GREEN
        pt.space_before = Pt(4)

    # Right Security & Isolation Card
    add_card(slide11, 6.7, 1.4, 6.0, 5.5, border_color=PRIMARY_PURPLE)
    tb_t2 = slide11.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.1))
    tf_t2 = tb_t2.text_frame
    tf_t2.word_wrap = True

    p = tf_t2.paragraphs[0]
    p.text = "Security & Isolation Standards"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_PURPLE

    sec_standards = [
        ("JWT Authentication Guard", "Protects all REST API endpoints. Unauthenticated requests are rejected with 401/403."),
        ("Bcrypt Password Encryption", "Passwords hashed with multi-round salt. Plaintext passwords never stored."),
        ("Strict User Data Isolation", "All SQL queries parameterized with user_id to prevent multi-tenant data leaks."),
        ("Cascading Foreign Keys", "Account deletion cascades cleanly across all 21 tables without orphaned records."),
        ("Environment Secret Protection", "API secrets and database keys isolated via .env configuration.")
    ]
    for title, desc in sec_standards:
        pt = tf_t2.add_paragraph()
        pt.text = "• " + title
        pt.font.size = Pt(13)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_WHITE
        pt.space_before = Pt(10)
        pd = tf_t2.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(11)
        pd.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 12: THANK YOU / Q&A (Closing Slide with Faded Hero Background)
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    apply_background(slide12)

    # Faded Background Image
    if os.path.exists(dash_img):
        bg_pic = slide12.shapes.add_picture(dash_img, Inches(1.5), Inches(0.8), Inches(10.333), Inches(5.8))

    # Center Hero Card
    card12 = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(1.2), Inches(9.333), Inches(5.1))
    card12.fill.solid()
    card12.fill.fore_color.rgb = BG_DARK
    card12.line.color.rgb = MINT_GREEN
    card12.line.width = Pt(2)

    tb12 = slide12.shapes.add_textbox(Inches(2.3), Inches(1.8), Inches(8.733), Inches(3.9))
    tf12 = tb12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "PennyWise"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_PURPLE

    p_sub = tf12.add_paragraph()
    p_sub.text = "Spend smarter. Think better. Save more."
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = MINT_GREEN
    p_sub.space_after = Pt(20)

    p_ty = tf12.add_paragraph()
    p_ty.text = "Thank You!"
    p_ty.alignment = PP_ALIGN.CENTER
    p_ty.font.size = Pt(36)
    p_ty.font.bold = True
    p_ty.font.color.rgb = TEXT_WHITE
    p_ty.space_after = Pt(10)

    p_qa = tf12.add_paragraph()
    p_qa.text = "Questions & Answers"
    p_qa.alignment = PP_ALIGN.CENTER
    p_qa.font.size = Pt(22)
    p_qa.font.bold = True
    p_qa.font.color.rgb = CYAN_BLUE

    print(f"[PPTX] Saving presentation to: {output_pptx}")
    prs.save(output_pptx)
    if os.path.exists(output_pptx):
        try:
            shutil.copyfile(output_pptx, target_pptx)
            print(f"[PPTX] Synced to: {target_pptx}")
        except Exception as e:
            print(f"[PPTX] Copy skipped: {e}")
    print("[PPTX] Presentation Generation Complete!")

if __name__ == '__main__':
    create_redesigned_presentation()
